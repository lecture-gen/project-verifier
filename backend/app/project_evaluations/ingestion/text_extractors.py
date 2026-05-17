from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from app.settings import ApiSettings


def looks_binary(data: bytes) -> bool:
    return b"\x00" in data[:2048]


def trim_text(value: str, settings: ApiSettings) -> str:
    return value[: settings.APP_MAX_EXTRACTED_TEXT_CHARS]


def extract_text(path: Path, settings: ApiSettings) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_text(path, settings)
    if suffix == ".docx":
        return extract_docx_text(path, settings)
    if suffix == ".pptx":
        return extract_pptx_text(path, settings)
    data = path.read_bytes()
    if looks_binary(data):
        return ""
    return trim_text(data.decode("utf-8", errors="ignore"), settings)


def extract_pdf_text(path: Path, settings: ApiSettings) -> str:
    reader = PdfReader(str(path))
    pages = []
    for index, page in enumerate(
        list(reader.pages)[: settings.APP_MAX_PDF_PAGES], start=1
    ):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[page {index}]\n{text}")
    return trim_text("\n\n".join(pages), settings)


def extract_docx_text(path: Path, settings: ApiSettings) -> str:
    document = Document(str(path))
    paragraphs = []
    for paragraph in document.paragraphs[: settings.APP_MAX_DOCX_PARAGRAPHS]:
        if paragraph.text:
            paragraphs.append(paragraph.text)
    return trim_text("\n".join(paragraphs), settings)


def extract_pptx_text(path: Path, settings: ApiSettings) -> str:
    presentation = Presentation(str(path))
    slides = []
    for slide_index, slide in enumerate(
        list(presentation.slides)[: settings.APP_MAX_PPTX_SLIDES], start=1
    ):
        texts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                texts.append(text)
        if texts:
            slides.append(f"[slide {slide_index}]\n" + "\n".join(texts))
    return trim_text("\n\n".join(slides), settings)
